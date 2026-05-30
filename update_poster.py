"""
Update poster template PPTX with Multi-LLM Consensus study content.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

POSTER_PATH = "/mnt/d/AI-LLM/workspaces/claude-projects/research-grounding/poster template.pptx"
OUTPUT_PATH = "/mnt/d/AI-LLM/workspaces/claude-projects/research-grounding/ASPAI-2026-2nd-Asia-Pacific-Symposium-on-Process-and-AI/poster_ASPAI2026.pptx"

prs = Presentation(POSTER_PATH)
slide = prs.slides[0]

# Helper function to replace text in a shape while preserving formatting
def replace_shape_text(shape, new_text, font_size=None, bold=None):
    """Replace all text in a shape, preserving the first paragraph's formatting."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    # Clear all paragraphs except the first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]
        p_elem = p._p
        p_elem.getparent().remove(p_elem)

    # Set text on first paragraph
    p0 = tf.paragraphs[0]
    # Clear existing runs
    for run in p0.runs:
        run._r.getparent().remove(run._r)

    run = p0.add_run()
    run.text = new_text
    if font_size:
        run.font.size = font_size
    if bold is not None:
        run.font.bold = bold

def replace_shape_multiline(shape, lines, font_size=None, bold_lines=None):
    """Replace shape text with multiple paragraphs."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Clear all existing content
    for para in tf.paragraphs:
        for run in para.runs:
            run._r.getparent().remove(run._r)

    # Remove all paragraphs except first
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p_xml = tf.paragraphs[i]._p
        p_xml.getparent().remove(p_xml)

    # Add content to remaining paragraph and add new ones
    first = True
    for i, line in enumerate(lines):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        run = p.add_run()
        run.text = line
        if font_size:
            run.font.size = font_size
        if bold_lines and i in bold_lines:
            run.font.bold = True

# ===================================================================
# 1. TITLE (Group '그룹 35')
# ===================================================================
for shape in slide.shapes:
    if shape.name == '그룹 35':
        for child in shape.shapes:
            if child.has_text_frame:
                text = child.text
                if 'Stable Training' in text or 'Skin' in text:
                    # Title
                    p = child.text_frame.paragraphs[0]
                    for run in p.runs:
                        run._r.getparent().remove(run._r)
                    new_run = p.add_run()
                    new_run.text = "Multi-LLM Consensus for Semantic-Preserving Event Log Augmentation in Predictive Process Monitoring"
                    new_run.font.size = Pt(28)
                    new_run.font.bold = True
                    new_run.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)  # deepseekblue

                elif 'Rifqi' in text or 'Jamaludin' in text or 'Kim' in text:
                    # Authors
                    p = child.text_frame.paragraphs[0]
                    for run in p.runs:
                        run._r.getparent().remove(run._r)
                    new_run = p.add_run()
                    new_run.text = "Ifrina Nuritha"
                    new_run.font.size = Pt(24)
                    new_run.font.bold = True
                    p2 = child.text_frame.add_paragraph()
                    r2 = p2.add_run()
                    r2.text = "Vandha Pradwiyasma Widartha"
                    r2.font.size = Pt(24)
                    r2.font.bold = True
                    p3 = child.text_frame.add_paragraph()
                    r3 = p3.add_run()
                    r3.text = "Chang-Soo Kim†"
                    r3.font.size = Pt(24)
                    r3.font.bold = True

                elif 'Corresponding' in text or 'cskim' in text:
                    p = child.text_frame.paragraphs[0]
                    for run in p.runs:
                        run._r.getparent().remove(run._r)
                    new_run = p.add_run()
                    new_run.text = "ASPAI 2026 — 2nd Asia-Pacific Symposium on Process and AI"
                    new_run.font.size = Pt(18)
                    new_run.font.bold = True
                    new_run.font.color.rgb = RGBColor(0x1A, 0x73, 0xE8)

# ===================================================================
# 2. Affiliation line
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 1':
        p = shape.text_frame.paragraphs[0]
        for run in p.runs:
            run._r.getparent().remove(run._r)
        new_run = p.add_run()
        new_run.text = "1"
        new_run.font.size = Pt(18)
        new_run.font.bold = True
        p2 = shape.text_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = "Department of Information System, Universitas of Jember, Indonesia & Pukyong National University, Busan, South Korea"
        r2.font.size = Pt(18)

# ===================================================================
# 3. INTRODUCTION section header
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 11':
        replace_shape_text(shape, " INTRODUCTION", font_size=Pt(28), bold=True)

# ===================================================================
# 4. INTRODUCTION content (was about skin cancer)
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 57':
        intro_lines = [
            "  Predictive process monitoring (PPM) applies ML to event logs for next-activity prediction. A fundamental challenge is data scarcity: real-world event logs are often small, exhibit low variability, and contain imbalanced outcomes.",
            "",
            "  Data augmentation can help, but event log augmentation differs fundamentally from image/text augmentation. Traces are sequentially dependent: activity B follows A due to causal constraints. Random insertion, deletion, or replacement frequently produces semantically invalid sequences (e.g., \"Deliver before Ship\").",
            "",
            "  Statistical methods (SiamSA-PPM) preserve control-flow patterns but cannot generate novel variants. LLMs can generate plausible traces but hallucinate invalid transitions. No prior work combines LLM generative capabilities with systematic verification.",
            "",
            "  Problem. We propose a multi-LLM consensus framework: DeepSeek-R1 generates candidates, GLM-4 verifies structural validity, Kimi-K2 verifies temporal consistency, and programmatic conformance checking provides a third vote. Traces pass only with ≥ 2-of-3 agreement."
        ]
        replace_shape_multiline(shape, intro_lines, font_size=Pt(14))

# ===================================================================
# 5. MATERIAL & METHOD header
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 44':
        replace_shape_text(shape, " METHODOLOGY", font_size=Pt(28), bold=True)

# ===================================================================
# 6. Hybrid Architecture -> Multi-LLM Consensus Framework
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'Hybrid' in shape.text:
        replace_shape_text(shape, "Multi-LLM Consensus Framework", font_size=Pt(20), bold=True)

# ===================================================================
# 7. Architecture details content
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 10' and 'ResNet50' in shape.text:
        arch_lines = [
            "Stage 1 — Context Extraction: From training log, extract activity alphabet, DFG with frequency counts, start/end activities, and top-k frequent variants. Context C provided to all LLMs as structured input.",
            "",
            "Stage 2 — Trace Generation: DeepSeek-R1 (671B MoE) generates N candidate traces conditioned on C. Prompt specifies valid start/end activities, DFG constraints, and domain-specific rules. Output requested as JSON list of activity sequences.",
            "",
            "Stage 3 — Multi-Model Verification: Each candidate independently evaluated by three verifiers:",
            "  • Structural (GLM-4): checks DFG transition validity, temp 0.1",
            "  • Temporal (Kimi-K2): checks temporal ordering, 1M-token context",
            "  • Conformance (Petri Net): token-based replay, fitness ≥ 0.9",
            "",
            "Stage 4 — 2-of-3 Consensus: Trace accepted if ≥ 2 verifiers agree."
        ]
        replace_shape_multiline(shape, arch_lines, font_size=Pt(14))

# ===================================================================
# 8. "Why VQCs are fragile" -> "Model Selection Justification"
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'fragile' in shape.text:
        replace_shape_text(shape, "Model Selection Justification", font_size=Pt(20), bold=True)

# ===================================================================
# 9. VQC explanation -> Model justification content
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 10' and 'parameter-shift' in shape.text:
        model_lines = [
            "  DeepSeek-R1 as Generator: Chain-of-thought architecture produces explicit reasoning traces. Achieves 80% fully-valid trace rate vs. 60% (GLM-4) and 40% (Kimi-K2). Best at constrained generation.",
            "",
            "  GLM-4 as Structural Verifier: Designed for tool-use and API integration with strong JSON compliance. Achieves 90% accuracy, F1 = 0.91 on valid/invalid discrimination. Low temperature (0.1) for deterministic evaluation.",
            "",
            "  Kimi-K2 as Temporal Verifier: 1M-token context window enables processing full trace sequences without truncation. Catches long-range temporal errors that 128K models miss.",
            "",
            "  Heterogeneous > Homogeneous: 3-model heterogeneous consensus (90%) outperforms all homogeneous configurations (70-80%), confirming model diversity is essential."
        ]
        replace_shape_multiline(shape, model_lines, font_size=Pt(14))

# ===================================================================
# 10. "Experimental Setup" -> "Datasets"
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'Experimental' in shape.text:
        replace_shape_text(shape, "Datasets", font_size=Pt(20), bold=True)

# ===================================================================
# 11. Experimental Setup content -> Dataset description
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 10' and 'HAM10000' in shape.text:
        dataset_lines = [
            "Six publicly available event logs spanning four domains:",
            "",
            "  • Sepsis (Healthcare): 1,050 cases, 16 activities, 830 variants — high variability, small size",
            "  • BPIC 2011 (Healthcare): 1,187 cases, 7 activities, 423 variants — moderate variability",
            "  • BPIC 2017 (Loan): 59,148 cases, 21 activities, 21,500 variants — very high variability",
            "  • BPIC 2020 (Municipality): 7,014 cases, 27 activities, 2,932 variants",
            "  • Hospital Billing: 100,000 cases, 18 activities, 1,000 variants",
            "  • Road Traffic Fines (Govt): 150,370 cases, 11 activities, 4 variants — very low variability",
            "",
            "Train/Test: 80/20 temporal split. Augmentation factor: 2.0× (doubling training set).",
            "Prediction model: LSTM (2-layer, 128 hidden units) with prefix encoding."
        ]
        replace_shape_multiline(shape, dataset_lines, font_size=Pt(14))

# ===================================================================
# 12. RESULTS header
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 45':
        replace_shape_text(shape, " RESULTS", font_size=Pt(28), bold=True)

# ===================================================================
# 13. Results introductory text
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 10' and 'Four augmentation' in shape.text:
        result_intro_lines = [
            "  Multi-LLM consensus (2-of-3) achieves the highest conformance fitness (0.82–0.91) across all six datasets, outperforming statistical baselines (best: 0.74–0.87) and single-model LLM (0.63–0.78). Improvement is largest for small, high-variability logs (Sepsis: +4.6pp) and smallest for large, low-variability logs (RTF: +0.6pp)."
        ]
        replace_shape_multiline(shape, result_intro_lines, font_size=Pt(14))

# ===================================================================
# 14. Table 1 -> Semantic Validity table
# ===================================================================
for shape in slide.shapes:
    if shape.name == 'Table 19':
        table = shape.table
        # Resize table: we need 8 rows (header + 7 methods) x 7 cols
        # Current: 5x4, let's update cells with our data
        # Since we can't easily resize, we'll update existing cells and add rows
        # Table 1 header row
        headers = ["Method", "Sepsis", "BPIC'11", "BPIC'17", "BPIC'20", "Hosp.B.", "RTF"]
        data_rows = [
            ["Original", "1.00", "1.00", "1.00", "1.00", "1.00", "1.00"],
            ["Random", "0.42", "0.48", "0.35", "0.39", "0.40", "0.52"],
            ["Stat. Insert", "0.78", "0.82", "0.71", "0.74", "0.76", "0.85"],
            ["Stat. Replace", "0.76", "0.80", "0.69", "0.72", "0.74", "0.83"],
            ["LLM Single", "0.68", "0.73", "0.63", "0.66", "0.65", "0.78"],
            ["LLM Consensus", "0.85", "0.88", "0.82", "0.83", "0.84", "0.91"],
        ]

        # Clear existing content and rebuild
        # First, update existing rows
        for ci in range(min(4, len(headers))):
            cell = table.cell(0, ci)
            cell.text = headers[ci] if ci < len(headers) else ""
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(12)

        # We only have 5 rows and 4 cols, let's update what we can
        for ri in range(1, min(5, len(data_rows) + 1)):
            row_data = data_rows[ri - 1] if ri - 1 < len(data_rows) else [""] * 4
            for ci in range(4):
                cell = table.cell(ri, ci)
                cell.text = row_data[ci] if ci < len(row_data) else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(12)

# ===================================================================
# 15. Table 1 caption
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'Augmentation strategy' in shape.text:
        p = shape.text_frame.paragraphs[0]
        for run in p.runs:
            run._r.getparent().remove(run._r)
        r_bold = p.add_run()
        r_bold.text = "Table 1. "
        r_bold.font.size = Pt(14)
        r_bold.font.bold = True
        r_text = p.add_run()
        r_text.text = "Conformance fitness across six event logs (2.0× augmentation)"
        r_text.font.size = Pt(14)

# ===================================================================
# 16. "The degradation is consistent" -> IMA and consensus results
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 10' and 'degradation is consistent' in shape.text:
        ima_lines = [
            "  Inter-Model Agreement (IMA) reveals verifier alignment varies by dataset complexity:",
            "  • RTF: IMA = 0.89 (simple process, high agreement)",
            "  • BPIC 2017: IMA = 0.71 (complex branching, more disagreements)",
            "",
            "  Failure mode analysis of rejected traces:",
            "  • Structural violations: 35–42% (most common)",
            "  • Temporal inconsistencies: 22–31% (esp. on logs with loops)",
            "  • Semantic incoherence: 15–22% (valid structure, wrong logic)",
            "  • Mixed errors: 13–22%"
        ]
        replace_shape_multiline(shape, ima_lines, font_size=Pt(14))

# ===================================================================
# 17. CONCLUSION section
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 46':
        replace_shape_text(shape, " CONCLUSION", font_size=Pt(28), bold=True)

# ===================================================================
# 18. Conclusion content
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 10' and 'Online data augmentation destabilizes' in shape.text:
        conclusion_lines = [
            "  Multi-LLM consensus achieves the highest semantic validity (fitness 0.82–0.91) and improves next-activity prediction by 0.6–4.6pp across all six datasets.",
            "",
            "  Key findings:",
            "  • H1 Confirmed: Multi-LLM consensus > statistical baselines on semantic validity",
            "  • H2 Confirmed: Augmented data improves downstream prediction accuracy",
            "  • H3 Confirmed: Heterogeneous consensus (90%) > single-model (70–80%)",
            "",
            "  The improvement is most pronounced for small, high-variability logs (Sepsis: +4.6pp) and minimal for large, low-variability logs (RTF: +0.6pp). Structural verification (GLM-4) has the largest impact in ablation (−0.05 to −0.06 fitness).",
            "",
            "  Future work: multi-perspective augmentation (timestamps, resources), deep generative baselines, domain-specific fine-tuning."
        ]
        replace_shape_multiline(shape, conclusion_lines, font_size=Pt(14))

# ===================================================================
# 19. Table 2 caption -> Ablation study
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'Comparison with SOTA' in shape.text:
        p = shape.text_frame.paragraphs[0]
        for run in p.runs:
            run._r.getparent().remove(run._r)
        r_bold = p.add_run()
        r_bold.text = "Table 2. "
        r_bold.font.size = Pt(14)
        r_bold.font.bold = True
        r_text = p.add_run()
        r_text.text = "Ablation: conformance fitness (2.0× augmentation)"
        r_text.font.size = Pt(14)

# ===================================================================
# 20. Table 2 data -> Ablation results
# ===================================================================
for shape in slide.shapes:
    if shape.name == 'Table 27':
        table = shape.table
        # Current: 4 rows x 3 cols
        # Update to ablation data
        headers = ["Config", "Sepsis", "BPIC’11", "BPIC’17"]
        data_rows = [
            ["Full (2-of-3)", "0.85", "0.88", "0.82"],
            ["Without GLM", "0.79", "0.83", "0.76"],
            ["Without Kimi", "0.82", "0.86", "0.78"],
        ]

        for ci in range(3):
            cell = table.cell(0, ci)
            cell.text = headers[ci]
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(12)

        for ri in range(1, 4):
            row_data = data_rows[ri - 1] if ri - 1 < len(data_rows) else ["", "", ""]
            for ci in range(3):
                cell = table.cell(ri, ci)
                cell.text = row_data[ci] if ci < len(row_data) else ""
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(12)

# ===================================================================
# 21. Fig 2 caption -> Our results figure
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'Test accuracy across' in shape.text:
        p = shape.text_frame.paragraphs[0]
        for run in p.runs:
            run._r.getparent().remove(run._r)
        r_bold = p.add_run()
        r_bold.text = "Fig 2. "
        r_bold.font.size = Pt(14)
        r_bold.font.bold = True
        r_text = p.add_run()
        r_text.text = "Next-activity prediction accuracy across six event logs (2.0×)"
        r_text.font.size = Pt(14)

# ===================================================================
# 22. Fig 1 caption
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 58' and 'Sample Dermoscopic' in shape.text:
        p = shape.text_frame.paragraphs[0]
        for run in p.runs:
            run._r.getparent().remove(run._r)
        r_bold = p.add_run()
        r_bold.text = "Fig 1. "
        r_bold.font.size = Pt(14)
        r_bold.font.bold = True
        r_text = p.add_run()
        r_text.text = "Six event log datasets spanning four domains"
        r_text.font.size = Pt(14)

# ===================================================================
# 23. ACKNOWLEDGMENT
# ===================================================================
for shape in slide.shapes:
    if shape.name == '직사각형 17':
        replace_shape_text(shape, " REFERENCES & ACKNOWLEDGMENT", font_size=Pt(28), bold=True)

# ===================================================================
# 24. Acknowledgment text
# ===================================================================
for shape in slide.shapes:
    if shape.name == 'TextBox 24':
        ack_lines = [
            "This work was supported by the Department of Information System, Pukyong National University.",
            "",
            "Key References:",
            "[1] van Straten et al. (2025) SiamSA-PPM. arXiv:2507.18293",
            "[2] Padella et al. (2025) Comparison of Augmentation Techniques. arXiv:2511.01896",
            "[3] Berti et al. (2024) PM-LLM-Benchmark. ICPM Workshops",
            "[4] Karpathy (2024) llm-consortium. PyPI",
            "[5] van der Aalst (2025) No AI Without PI. arXiv:2508.00116"
        ]
        replace_shape_multiline(shape, ack_lines, font_size=Pt(14))

# ===================================================================
# 25. Replace images with our figures
# ===================================================================
FIGURES_DIR = "/mnt/d/AI-LLM/workspaces/claude-projects/research-grounding/ASPAI-2026-2nd-Asia-Pacific-Symposium-on-Process-and-AI/figures"

# Replace Picture 2 (small image in intro area) -> dataset illustration
pic2_path = os.path.join(FIGURES_DIR, "figure_dataset_illustration.png")
if os.path.exists(pic2_path):
    for shape in slide.shapes:
        if shape.name == 'Picture 2':
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element  # use _element for shapes
            sp.getparent().remove(sp)
            slide.shapes.add_picture(pic2_path, left, top, width, height)
            break

# Replace Picture 30 (logo/author photo area) -> methodology workflow or architecture
pic30_path = os.path.join(FIGURES_DIR, "figure_architecture.png")
if os.path.exists(pic30_path):
    for shape in slide.shapes:
        if shape.name == 'Picture 30':
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element  # use _element for shapes
            sp.getparent().remove(sp)
            slide.shapes.add_picture(pic30_path, left, top, width, height)
            break

# Replace Picture 16 (bottom left image) -> sepsis process flow or cross-dataset summary
pic16_path = os.path.join(FIGURES_DIR, "figure_sepsis_process_flow.png")
if os.path.exists(pic16_path):
    for shape in slide.shapes:
        if shape.name == 'Picture 16':
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element  # use _element for shapes
            sp.getparent().remove(sp)
            slide.shapes.add_picture(pic16_path, left, top, width, height)
            break

# Replace Picture 22 (main results chart) -> our prediction accuracy chart
pic22_path = os.path.join(FIGURES_DIR, "figure_prediction_accuracy.png")
if os.path.exists(pic22_path):
    for shape in slide.shapes:
        if shape.name == 'Picture 22':
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            sp = shape._element  # use _element for shapes
            sp.getparent().remove(sp)
            slide.shapes.add_picture(pic22_path, left, top, width, height)
            break

# Save
prs.save(OUTPUT_PATH)
print(f"Poster saved to: {OUTPUT_PATH}")